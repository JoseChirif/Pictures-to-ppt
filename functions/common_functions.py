import os
from pptx import Presentation
from pptx.util import Cm, Pt, Inches
from pptx.dml.color import RGBColor
from PIL import Image

def split_filename(filename):
    parts = filename.split('_', 1)
    if len(parts) > 1:
        return parts[0], parts[1]
    else:
        return parts[0], ''

def setup_presentation(width_cm, height_cm):
    prs = Presentation()
    prs.slide_width = Cm(width_cm)
    prs.slide_height = Cm(height_cm)
    return prs

def add_image_to_slide(slide, image_path, prs, fixed_height=None, border=False):
    with Image.open(image_path) as img:
        img_width, img_height = img.size

    if fixed_height is None:
        new_height = prs.slide_height
        new_width = int((img_width / img_height) * new_height)
    else:
        new_height = Cm(fixed_height)
        new_width = int((img_width / img_height) * new_height)

    if new_width > prs.slide_width:
        new_width = prs.slide_width
        new_height = int((img_height / img_width) * new_width)

    left = int((prs.slide_width - new_width) / 2)
    top = int((prs.slide_height - new_height) / 2)

    picture = slide.shapes.add_picture(image_path, left, top, width=new_width, height=new_height)

    if border:
        line = picture.line
        line.color.rgb = RGBColor(0, 0, 0)
        line.width = Pt(0.75)

def add_text_to_slide(slide, text, font_size, color, bg_color, prs, position='top', alignment='center'):
    if position == 'top':
        text_left = Cm(1)
        text_top = Cm(1)
    elif position == 'bottom':
        text_left = Cm(0)
        text_top = prs.slide_height - Cm(1.5)

    text_width = prs.slide_width
    text_height = Cm(1.5)

    txBox = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
    tf = txBox.text_frame
    tf.clear()

    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.name = 'Aptos'

    fill = txBox.fill
    fill.solid()
    fill.fore_color.rgb = bg_color

    if alignment == 'center':
        p.alignment = 1
    elif alignment == 'right':
        p.alignment = 2
