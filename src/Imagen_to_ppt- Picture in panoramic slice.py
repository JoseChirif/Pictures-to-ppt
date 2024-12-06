import os
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from PIL import Image

# go to the parent directory if you are running this script directly (uncomment the following lines)
import sys
sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from config.config import pictures_extensions
from functions.functions import working_directory, show_options, show_message, check_presentation_exists


def create_presentation():
    """
    Creates a PowerPoint presentation from images in the working directory.

    Args:
        ppt_name (str): Name of the PowerPoint presentation (without extension).
    """
    
    # Variables
    ppt_name = "Presentation - pictures in panoramic slides"
    
    # Get the working directory and construct the output file path
    image_folder = working_directory()
    output_pptx = os.path.join(image_folder, ppt_name + ".pptx")
    
    # Check if the file exists and handle replacement decision
    if not check_presentation_exists(output_pptx):
        show_message("Operation canceled", "The presentation was not created.")
        return

    
    # Create a blank presentation with standard dimensions
    prs = Presentation()
    prs.slide_width = Cm(33.867)
    prs.slide_height = Cm(19.05)
    
    # Get the list of images in the folder
    images = [img for img in os.listdir(image_folder) if img.lower().endswith(pictures_extensions)]
    
    # Get option nr
    option_nr = show_options("Do you want to include \n the file name's in the slide?", "Yes", "No")
    # If `option_nr` is None, cancel the operation
    if option_nr is None:
        show_message("Operation canceled", "No presentation was created.")
        return
    
    for image in images:
        # Full path of the image
        image_path = os.path.join(image_folder, image)
        
        # Get image dimensions
        with Image.open(image_path) as img:
            img_width, img_height = img.size
        
        # Calculate the new width keeping the aspect ratio
        new_height = prs.slide_height
        new_width = int((img_width / img_height) * new_height)
        
        if new_width > prs.slide_width:
            new_width = prs.slide_width
            new_height = int((img_height / img_width) * new_width)
        
        # Calculate the positions to center the image
        left = int((prs.slide_width - new_width) / 2)
        top = int((prs.slide_height - new_height) / 2)
        
        # Create a new slide
        slide_layout = prs.slide_layouts[6]  # Use a blank slide layout (no title or text)
        slide = prs.slides.add_slide(slide_layout)
        
        # Add the image to the slide
        slide.shapes.add_picture(image_path, left, top, width=new_width, height=new_height)
        
        
        if option_nr == 1:
            # Obtener el nombre del archivo sin la extensión
            file_name = os.path.splitext(image)[0]  
        
            
            # Añadir el cuadro de texto (parte superior)
            text_left = Cm(1.02)  # 1.02 cm desde la izquierda
            text_top = Cm(1.14)   # 1.14 cm desde el borde superior
            text_width = Cm(5)  # Ancho del cuadro de texto
            text_height = Cm(2)  # Alto del cuadro de texto
            
            txBox = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
            tf = txBox.text_frame
            tf.clear()  # Asegurarse de que el cuadro de texto esté vacío
            
            p = tf.add_paragraph()
            p.text = file_name
            p.font.size = Pt(28)  # Tamaño del texto
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 0, 0)  # Texto negro
            p.font.name = 'Aptos'  # Fuente Aptos
            
            # Añadir fondo naranja al cuadro de texto
            fill = txBox.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(250, 192, 144)  # Color de fondo #FAC090
            
            # Alinear el texto en el cuadro de texto
            txBox.text_frame.paragraphs[0].alignment = 1  # Centrado
            
        elif option_nr == 2:
            continue
        
    # Guardar la presentación
    prs.save(output_pptx)
    show_message("Done", f"Presentation saved as {ppt_name + '.pptx'}")


if __name__ == "__main__":
    create_presentation()


