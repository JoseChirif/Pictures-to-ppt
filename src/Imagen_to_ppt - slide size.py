import os
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from PIL import Image

def create_presentation(image_folder, output_pptx):
    prs = Presentation()
    
    images = [img for img in os.listdir(image_folder) if img.lower().endswith('.jpg')]
    
    max_width_in_inches = 56  # Máximo tamaño permitido en pulgadas
    max_height_in_inches = 56
    
    for image in images:
        image_path = os.path.join(image_folder, image)
        
        with Image.open(image_path) as img:
            width, height = img.size
            dpi = img.info.get('dpi', (72, 72))
            width_in_inches = width / dpi[0]
            height_in_inches = height / dpi[1]
            
            # Limitar el tamaño de la diapositiva
            if width_in_inches > max_width_in_inches:
                width_in_inches = max_width_in_inches
            if height_in_inches > max_height_in_inches:
                height_in_inches = max_height_in_inches
        
        prs.slide_width = Inches(width_in_inches)
        prs.slide_height = Inches(height_in_inches)
        
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        # Ajustar el tamaño de la imagen al tamaño de la diapositiva
        slide.shapes.add_picture(image_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
        
        file_name = os.path.splitext(image)[0]
        left = Cm(1)
        top = Cm(1)
        width = Inches(2)
        height = Inches(0.5)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = file_name
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 0, 0)
        
        fill = txBox.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(250, 192, 144)
        
        txBox.text_frame.paragraphs[0].alignment = 1
        
    prs.save(output_pptx)
    print(f"Presentación guardada como {output_pptx}")

if __name__ == "__main__":
    image_folder = "imagenes"
    output_pptx = "presentacion.pptx"
    create_presentation(image_folder, output_pptx)