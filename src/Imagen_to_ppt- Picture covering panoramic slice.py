import os
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from PIL import Image

# go to the parent directory if you are running this script directly (uncomment the following lines)
# import sys
# sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from config.config import pictures_extensions
from functions.functions import working_directory, show_options, show_message, check_presentation_exists


def create_presentation():
    """
    Creates a PowerPoint presentation from images in the working directory.

    Args:
        ppt_name (str): Name of the PowerPoint presentation (without extension).
    """
    
    # Variables
    ppt_name = "Presentation - pictures covering panoramic slides"
    
    # Get the working directory and construct the output file path
    image_folder = working_directory()
    output_pptx = os.path.join(image_folder, ppt_name + ".pptx")
    
    # Check if the file exists and handle replacement decision
    if not check_presentation_exists(output_pptx):
        show_message("Operation canceled", "The presentation was not created.")
        return
    
    # Create a blank presentation with standard dimensions
    prs = Presentation()
    prs.slide_width = Cm(33.867)  # 16:9 aspect ratio
    prs.slide_height = Cm(19.05)
    
    # Get the list of images in the folder
    images = [img for img in os.listdir(image_folder) if img.lower().endswith(pictures_extensions)]
    
    # Get option nr
    option_nr = show_options("Do you want to include \n the file name's in the slide?", "Yes", "No")
    if option_nr is None:
        show_message("Operation canceled", "No presentation was created.")
        return
    
    for image in images:
        # Full path of the image
        image_path = os.path.join(image_folder, image)
        
        # Get image dimensions
        with Image.open(image_path) as img:
            img_width, img_height = img.size
        
        # Slide dimensions
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        slide_aspect_ratio = slide_width / slide_height
        
        # Image aspect ratio
        img_aspect_ratio = img_width / img_height
        
        # Adjust the image to cover the slide
        if img_aspect_ratio > slide_aspect_ratio:
            # Image is wider: fit height, crop sides
            new_height = slide_height
            new_width = (img_width / img_height) * new_height
            left = int((slide_width - new_width) / 2)
            top = 0
        else:
            # Image is taller: fit width, crop top/bottom
            new_width = slide_width
            new_height = (img_height / img_width) * new_width
            top = int((slide_height - new_height) / 2)
            left = 0
        
        # Create a new slide
        slide_layout = prs.slide_layouts[6]  # Use a blank slide layout (no title or text)
        slide = prs.slides.add_slide(slide_layout)
        
        # Add the image to the slide, cropped to cover the entire slide
        slide.shapes.add_picture(image_path, left, top, width=new_width, height=new_height)
        
        if option_nr == 1:
            # Get the file name without extension
            file_name = os.path.splitext(image)[0]
            
            # Add a text box at the top
            text_left = Cm(1.02)
            text_top = Cm(1.14)
            text_width = Cm(5)
            text_height = Cm(2)
            
            txBox = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
            tf = txBox.text_frame
            tf.clear()
            
            p = tf.add_paragraph()
            p.text = file_name
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 0, 0)
            p.font.name = 'Aptos'
            
            # Add orange background to the text box
            fill = txBox.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(250, 192, 144)  # Color #FAC090
            
            # Center-align the text
            txBox.text_frame.paragraphs[0].alignment = 1
        
    # Save the presentation
    prs.save(output_pptx)
    show_message("Done", f"Presentation saved as {ppt_name + '.pptx'}")





if __name__ == "__main__":
    create_presentation()

