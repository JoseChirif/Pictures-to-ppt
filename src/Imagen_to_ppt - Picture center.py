import os
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from PIL import Image

# go to the parent directory if you are running this script directly (uncomment the following lines)
# import sys
# sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from config.config import pictures_extensions
from functions.functions import working_directory, show_options, show_message, check_presentation_exists



def create_presentation():
    """
    Creates a PowerPoint presentation from images in the working directory.

    Args:
        ppt_name (str): Name of the PowerPoint presentation (without extension).
    """
    # Variables
    ppt_name = "Presentation - Pictures center ppt"
    
    # Get the working directory and construct the output file path
    image_folder = working_directory()
    output_pptx = os.path.join(image_folder, ppt_name + ".pptx")
    
    # Check if the file exists and handle replacement decision
    if not check_presentation_exists(output_pptx):
        show_message("Operation canceled", "The presentation was not created.")
        return

    # Create a blank presentation with standard dimensions
    prs = Presentation()
    prs.slide_width = Cm(33.867)
    prs.slide_height = Cm(19.05)
    
    # Get the list of images in the folder
    images = [img for img in os.listdir(image_folder) if img.lower().endswith(pictures_extensions)]
    
    # Get option nr
    option_nr = show_options("Do you want to include \n the file name's in the slide?", "Yes", "No")
    # If `option_nr` is None, cancel the operation
    if option_nr is None:
        show_message("Operation canceled", "No presentation was created.")
        return
    
    for image in images:
        # Full path of the image
        image_path = os.path.join(image_folder, image)
        
        # Get image dimensions
        with Image.open(image_path) as img:
            img_width, img_height = img.size
        
        # Set the fixed height for the image
        new_height = Cm(17.43)
        new_width = int((img_width / img_height) * new_height)
        
        if new_width > prs.slide_width:
            new_width = prs.slide_width
            new_height = int((img_height / img_width) * new_width)
        
        # Calculate positions to center the image
        left = int((prs.slide_width - new_width) / 2)
        top = int((prs.slide_height - new_height) / 2)
        
        # Create a new slide
        slide_layout = prs.slide_layouts[6]  # Use a blank slide layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Add the centered image with a black border
        picture = slide.shapes.add_picture(image_path, left, top, width=new_width, height=new_height)
        line = picture.line
        line.color.rgb = RGBColor(0, 0, 0)  # Black border
        line.width = Pt(0.75)  # Border thickness 0.75 pt
        
        if option_nr == 1:
            # Get the file name without the extension
            file_name = os.path.splitext(image)[0]
            
            # Add the file name at the bottom of the slide
            text_left = Cm(0)  # Horizontally center the text
            text_top = prs.slide_height - Cm(1.5)  # Position at the bottom
            text_width = prs.slide_width
            text_height = Cm(1)
            
            txBox = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
            tf = txBox.text_frame
            tf.clear()
            
            p = tf.add_paragraph()
            p.text = file_name
            p.font.size = Pt(10)  # Text size 10
            p.font.color.rgb = RGBColor(0, 0, 0)  # Black text
            p.font.name = 'Aptos'  # Aptos font
            p.alignment = 2  # Centered
        
        elif option_nr == 2:
            continue
    
    # Save the presentation
    prs.save(output_pptx)
    show_message("Done", f"Presentation saved as {ppt_name + '.pptx'}")



if __name__ == "__main__":
    create_presentation()
