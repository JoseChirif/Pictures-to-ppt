import os
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from PIL import Image

def split_filename(filename):
    parts = filename.split('_', 1)
    if len(parts) > 1:
        return parts[0], parts[1]
    else:
        return parts[0], ''

def create_presentation(image_folder, output_pptx):
    # Crear una presentación en blanco con dimensiones específicas
    prs = Presentation()
    prs.slide_width = Cm(33.867)
    prs.slide_height = Cm(19.05)
    
    # Obtener la lista de imágenes en la carpeta
    images = [img for img in os.listdir(image_folder) if img.lower().endswith('.jpg')]
    
    for image in images:
        # Ruta completa de la imagen
        image_path = os.path.join(image_folder, image)
        
        # Obtener dimensiones de la imagen
        with Image.open(image_path) as img:
            img_width, img_height = img.size
        
        # Calcular el nuevo ancho manteniendo la relación de aspecto
        new_height = prs.slide_height
        new_width = int((img_width / img_height) * new_height)
        
        if new_width > prs.slide_width:
            new_width = prs.slide_width
            new_height = int((img_height / img_width) * new_width)
        
        # Calcular las posiciones para centrar la imagen
        left = int((prs.slide_width - new_width) / 2)
        top = int((prs.slide_height - new_height) / 2)
        
        # Crear una nueva diapositiva
        slide_layout = prs.slide_layouts[6]  # Usar un diseño de diapositiva en blanco (sin título ni texto)
        slide = prs.slides.add_slide(slide_layout)
        
        # Añadir la imagen a la diapositiva
        slide.shapes.add_picture(image_path, left, top, width=new_width, height=new_height)
        
        # Obtener el nombre del archivo sin la extensión
        file_name = os.path.splitext(image)[0]  
      
        
        # Añadir el cuadro de texto (parte superior)
        text_left = Cm(1.02)  # 1.02 cm desde la izquierda
        text_top = Cm(1.14)   # 1.14 cm desde el borde superior
        text_width = Cm(5)  # Ancho del cuadro de texto
        text_height = Cm(2)  # Alto del cuadro de texto
        
        txBox = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
        tf = txBox.text_frame
        tf.clear()  # Asegurarse de que el cuadro de texto esté vacío
        
        p = tf.add_paragraph()
        p.text = file_name
        p.font.size = Pt(28)  # Tamaño del texto
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 0, 0)  # Texto negro
        p.font.name = 'Aptos'  # Fuente Aptos
        
        # Añadir fondo naranja al cuadro de texto
        fill = txBox.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(250, 192, 144)  # Color de fondo #FAC090
        
        # Alinear el texto en el cuadro de texto
        txBox.text_frame.paragraphs[0].alignment = 1  # Centrado
        
      
        
    # Guardar la presentación
    prs.save(output_pptx)
    print(f"Presentación guardada como {output_pptx}")

if __name__ == "__main__":
    image_folder = "imagenes"  # Carpeta de imágenes en el mismo directorio que el script
    output_pptx = "presentacion.pptx"  # Nombre del archivo de salida
    create_presentation(image_folder, output_pptx)
